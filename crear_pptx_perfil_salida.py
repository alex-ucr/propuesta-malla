"""Genera PPTX con las tres tablas de la subsección Perfil de salida."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUTPUT = "perfil-salida.pptx"

HEADER_COLOR = RGBColor(0x00, 0xC0, 0xF3)
HEADER_TEXT = RGBColor(0x00, 0x00, 0x00)
ROW_ALT = RGBColor(0xE8, 0xF7, 0xFD)
ROW_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0x99, 0x99, 0x99)

TABLES = [
    {
        "title": "Dimensión declarativa",
        "subtitle": "Tabla de conocimientos declarativos",
        "col_headers": ("Código(s)", "Enunciado de la dimensión declarativa"),
        "rows": [
            ("CD01", "Conoce y comprende los fundamentos matemáticos —incluyendo la lógica, la teoría de números y la geometría euclídea— y los principios formales que sustentan el razonamiento y la argumentación rigurosa."),
            ("CD02", "Conoce el uso de software especializado, los fundamentos de la programación y la computación científica, así como las técnicas básicas de análisis numérico para formular, explorar y resolver problemas matemáticos mediante enfoques computacionales."),
            ("CD03", "Conoce los conceptos, métodos y aplicaciones básicos del cálculo y del álgebra lineal como fundamentos de la formación matemática."),
            ("CD04", "Conoce los conceptos y técnicas de las ecuaciones diferenciales, la teoría de la medida, la probabilidad y los procesos estocásticos, así como del análisis complejo y funcional."),
            ("CD05", "Conoce los fundamentos y resultados básicos de la geometría y la topología."),
            ("CD06", "Conoce las estructuras y resultados fundamentales del álgebra abstracta."),
            ("CD07", "Conoce los principios de la matemática aplicada, por ejemplo por medio de la programación, la modelación, computación científica o el análisis de datos, para resolver problemas y tomar decisiones en ámbitos como la industria, la banca y otros sectores."),
        ],
    },
    {
        "title": "Dimensión procedimental",
        "subtitle": "Tabla de conocimientos procedimentales",
        "col_headers": ("Código(s)", "Enunciado de la dimensión procedimental"),
        "rows": [
            ("CP01", "Aplica pensamiento operacional, abstracto y formal para analizar, resolver y sintetizar problemas matemáticos."),
            ("CP02", "Formula conjeturas, resuelve problemas, interpreta resultados y se comunica de manera coherente en lenguaje matemático formal."),
            ("CP03", "Utiliza tecnologías y software científico para explorar, validar y resolver problemas matemáticos."),
            ("CP04", "Participa activamente en discusiones científicas y divulga el conocimiento matemático de forma accesible a diferentes audiencias."),
        ],
    },
    {
        "title": "Dimensión actitudinal",
        "subtitle": "Tabla de conocimientos actitudinales",
        "col_headers": ("Código(s)", "Enunciado de la dimensión actitudinal"),
        "rows": [
            ("CA01", "Manifiesta respeto por la diversidad humana, matemática, científica y cultural, fomentando un ambiente inclusivo y de valoración mutua."),
            ("CA02", "Se comunica con claridad, colabora en equipos académicos y profesionales, y mantiene una disposición al autoaprendizaje y al liderazgo."),
            ("CA03", "Demuestra conciencia social y participación ciudadana, con una visión sociocultural e integral de la sociedad en la que la ciencia y la matemática ocupan un lugar relevante."),
            ("CA04", "Ejerce pensamiento crítico, actúa con integridad ética y evalúa con rigor las prácticas científicas y académicas."),
            ("CA05", "Utiliza su conocimiento matemático para resolver problemas académicos y profesionales, adaptándolo a distintos contextos y necesidades."),
            ("CA06", "Mantiene un interés permanente por los desarrollos de la investigación científica y matemática, valorando la actualización como parte de su formación profesional."),
        ],
    },
]


def set_cell_fill(cell, color):
    cell.fill.solid()
    cell.fill.fore_color.rgb = color


def set_cell_text(cell, text, *, bold=False, size=11, align=PP_ALIGN.LEFT):
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = "Calibri"
    run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Perfil de salida"
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.name = "Calibri"

    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(8.4), Inches(1.5))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = (
        "Síntesis del perfil de egreso (2018–2020)\n"
        "Bachillerato y Licenciatura en Matemática\n"
        "Escuela de Matemática — UCR"
    )
    run.font.size = Pt(18)
    run.font.name = "Calibri"


def add_table_slide(prs, table_info):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Slide title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = table_info["title"]
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.name = "Calibri"

    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.75), Inches(9), Inches(0.35))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = table_info["subtitle"]
    run.font.size = Pt(14)
    run.font.italic = True
    run.font.name = "Calibri"

    rows_data = table_info["rows"]
    n_rows = len(rows_data) + 1
    n_cols = 2

    left = Inches(0.4)
    top = Inches(1.15)
    width = Inches(9.2)
    height = Inches(5.9)

    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = shape.table

    table.columns[0].width = Inches(1.0)
    table.columns[1].width = Inches(8.2)

    # Header row
    for col_idx, header_text in enumerate(table_info["col_headers"]):
        cell = table.cell(0, col_idx)
        set_cell_fill(cell, HEADER_COLOR)
        align = PP_ALIGN.CENTER if col_idx == 0 else PP_ALIGN.LEFT
        set_cell_text(cell, header_text, bold=True, size=11, align=align)

    # Data rows
    for row_idx, (code, statement) in enumerate(rows_data, start=1):
        bg = ROW_ALT if row_idx % 2 == 0 else ROW_WHITE
        code_cell = table.cell(row_idx, 0)
        stmt_cell = table.cell(row_idx, 1)
        set_cell_fill(code_cell, bg)
        set_cell_fill(stmt_cell, bg)
        set_cell_text(code_cell, code, bold=True, size=11, align=PP_ALIGN.CENTER)
        set_cell_text(stmt_cell, statement, size=10, align=PP_ALIGN.LEFT)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    for table_info in TABLES:
        add_table_slide(prs, table_info)

    prs.save(OUTPUT)
    print(f"Guardado: {OUTPUT}")


if __name__ == "__main__":
    main()
